

import os
import base64
import zipfile

from unstructured.partition.html import partition_html
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.docx import partition_docx
from unstructured.partition.pptx import partition_pptx
from unstructured.partition.text import partition_text
from unstructured.partition.md import partition_md
from unstructured.documents.elements import Image as UnstructuredImage, ElementMetadata

from src.services.llm import openAI
from src.config.logging import get_logger
from langchain_core.messages import HumanMessage

logger = get_logger(__name__)

_VALID_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}


# ─────────────────────────────────────────────
#  PPTX helpers
# ─────────────────────────────────────────────

def _iter_pptx_shapes(shapes):
    """Recursively yield shapes, descending into group shapes."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _extract_pptx_images_with_slide_numbers(file_path: str) -> list:
    """
    Dual-pass PPTX image extraction.

    Pass 1 (python-pptx): catches PICTURE shapes, PLACEHOLDER shapes that
    hold images, and shapes with a fill image.  Records each image's slide
    number so it can be placed into the right chunk later.

    Pass 2 (ZIP fallback): scans ppt/media/ and adds any image file that was
    NOT already captured in Pass 1 (deduped by raw bytes hash).
    """
    logger.info("pptx_image_extraction_started", file_path=file_path)

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("python_pptx_not_installed_falling_back_to_zip", file_path=file_path)
        return _extract_images_from_office_zip(file_path, "pptx")

    image_elements: list = []
    seen_hashes: set = set()

    def _add_image(image_bytes: bytes, slide_num: int, label: str = ""):
        if not image_bytes:
            return
        h = hash(image_bytes)
        if h in seen_hashes:
            return
        seen_hashes.add(h)
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        metadata = ElementMetadata(image_base64=b64, page_number=slide_num)
        image_elements.append(UnstructuredImage(text="", metadata=metadata))
        logger.debug(
            "pptx_image_extracted",
            slide=slide_num,
            bytes=len(image_bytes),
            label=label,
        )

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        logger.info("pptx_opened", total_slides=total_slides)

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            for shape in _iter_pptx_shapes(slide.shapes):
                try:
                    stype = shape.shape_type

                    # Case 1: explicit picture shape
                    if stype == MSO_SHAPE_TYPE.PICTURE:
                        _add_image(shape.image.blob, slide_num, label="PICTURE")
                        continue

                    # Case 2: placeholder that contains an image (e.g. content placeholder)
                    if stype == MSO_SHAPE_TYPE.PLACEHOLDER:
                        try:
                            _add_image(shape.image.blob, slide_num, label="PLACEHOLDER")
                        except Exception:
                            pass  # placeholder doesn't hold an image — skip silently
                        continue

                    # Case 3: shape fill is a picture (background images, watermarks, etc.)
                    try:
                        fill = shape.fill
                        if fill and hasattr(fill, "type") and str(fill.type) == "PICTURE (7)":
                            _add_image(fill.fore_color._element.blob, slide_num, label="FILL")
                    except Exception:
                        pass

                except Exception as shape_err:
                    logger.debug(
                        "pptx_shape_skipped",
                        slide=slide_num,
                        error=str(shape_err),
                    )

        logger.info("pptx_pass1_done", images_found=len(image_elements))

        # Pass 2: ZIP sweep — pick up anything python-pptx missed
        zip_elements = _extract_images_from_office_zip(file_path, "pptx", seen_hashes=seen_hashes)
        if zip_elements:
            logger.info("pptx_pass2_recovered_additional_images", count=len(zip_elements))
            image_elements.extend(zip_elements)

        logger.info("pptx_image_extraction_done", total_images=len(image_elements))
        return image_elements

    except Exception as e:
        logger.error(
            "pptx_image_extraction_failed_falling_back",
            error=str(e),
            exc_info=True,
        )
        return _extract_images_from_office_zip(file_path, "pptx")


# ─────────────────────────────────────────────
#  DOCX helpers
# ─────────────────────────────────────────────

def _extract_docx_images_with_python_docx(file_path: str) -> list:
    """Extract images from a .docx file via python-docx relationship scanning."""
    logger.info("docx_image_extraction_started", file_path=file_path)

    try:
        from docx import Document
    except ImportError:
        logger.warning("python_docx_not_installed_falling_back_to_zip", file_path=file_path)
        return _extract_images_from_office_zip(file_path, "docx")

    image_elements: list = []

    try:
        doc = Document(file_path)
        for rel_id, rel in doc.part.rels.items():
            if "image" not in rel.reltype.lower():
                continue
            try:
                image_part = rel.target_part
                image_bytes = image_part.blob
                if not image_bytes:
                    continue
                ext = os.path.splitext(getattr(image_part, "partname", ""))[1].lower()
                if ext and ext not in _VALID_IMAGE_EXTENSIONS:
                    logger.debug("docx_skipping_non_image_rel", rel_id=rel_id, ext=ext)
                    continue
                b64 = base64.b64encode(image_bytes).decode("utf-8")
                metadata = ElementMetadata(image_base64=b64)
                image_elements.append(UnstructuredImage(text="", metadata=metadata))
                logger.debug("docx_image_extracted", rel_id=rel_id, bytes=len(image_bytes))
            except Exception as rel_err:
                logger.debug("docx_rel_skipped", rel_id=rel_id, error=str(rel_err))

        logger.info("docx_image_extraction_done", total_images=len(image_elements))
        return image_elements

    except Exception as e:
        logger.error(
            "docx_image_extraction_failed_falling_back",
            error=str(e),
            exc_info=True,
        )
        return _extract_images_from_office_zip(file_path, "docx")


# ─────────────────────────────────────────────
#  ZIP fallback (shared by docx + pptx)
# ─────────────────────────────────────────────

def _extract_images_from_office_zip(
    file_path: str, file_type: str, seen_hashes: set = None
) -> list:
    """
    Fallback: treat the Office file as a ZIP and pull images from the media folder.
    `seen_hashes` lets the caller skip images already captured in an earlier pass.
    """
    media_folder = (
        "word/media/" if file_type == "docx"
        else "ppt/media/" if file_type == "pptx"
        else None
    )
    if not media_folder:
        return []

    if seen_hashes is None:
        seen_hashes = set()

    image_elements: list = []
    logger.info("zip_fallback_started", file_type=file_type, media_folder=media_folder)

    try:
        with zipfile.ZipFile(file_path, "r") as archive:
            for name in archive.namelist():
                if not name.startswith(media_folder):
                    continue
                ext = os.path.splitext(name)[1].lower()
                if ext not in _VALID_IMAGE_EXTENSIONS:
                    continue
                try:
                    with archive.open(name) as f:
                        image_bytes = f.read()
                    if not image_bytes:
                        continue
                    h = hash(image_bytes)
                    if h in seen_hashes:
                        continue
                    seen_hashes.add(h)
                    b64 = base64.b64encode(image_bytes).decode("utf-8")
                    metadata = ElementMetadata(
                        image_base64=b64, filename=os.path.basename(name)
                    )
                    image_elements.append(UnstructuredImage(text="", metadata=metadata))
                    logger.debug("zip_image_extracted", name=name, bytes=len(image_bytes))
                except Exception as inner_err:
                    logger.debug("zip_entry_skipped", name=name, error=str(inner_err))

        logger.info("zip_fallback_done", images_found=len(image_elements))
    except Exception as e:
        logger.error("zip_fallback_failed", error=str(e), exc_info=True)

    return image_elements


# ─────────────────────────────────────────────
#  Core partition entry point
# ─────────────────────────────────────────────

def partition_document(temp_file: str, file_type: str, source_type: str = "file") -> list:
    """
    Partition a document into unstructured elements.
    For docx/pptx, image extraction is performed separately and merged in,
    because the unstructured library does not natively extract embedded images
    from those formats.
    """
    source = (source_type or "file").lower()
    kind = (file_type or "").lower()

    logger.info(
        "partition_document_started",
        file=temp_file,
        file_type=kind,
        source_type=source,
    )

    if source == "url":
        return partition_html(filename=temp_file)

    if kind == "pdf":
        elements = partition_pdf(
            filename=temp_file,
            strategy="hi_res",
            infer_table_structure=True,
            extract_image_block_types=["Image"],
            extract_image_block_to_payload=True,
        )
        logger.info("partition_pdf_done", element_count=len(elements))
        return elements

    elif kind == "docx":
        elements = partition_docx(
            filename=temp_file, strategy="hi_res", infer_table_structure=True
        )
        image_elements = _extract_docx_images_with_python_docx(temp_file)
        elements.extend(image_elements)
        logger.info(
            "partition_docx_done",
            text_table_elements=len(elements) - len(image_elements),
            image_elements=len(image_elements),
            total=len(elements),
        )
        return elements

    elif kind == "pptx":
        elements = partition_pptx(
            filename=temp_file, strategy="hi_res", infer_table_structure=True
        )
        image_elements = _extract_pptx_images_with_slide_numbers(temp_file)
        elements.extend(image_elements)
        logger.info(
            "partition_pptx_done",
            text_table_elements=len(elements) - len(image_elements),
            image_elements=len(image_elements),
            total=len(elements),
        )
        return elements

    elif kind == "txt":
        return partition_text(filename=temp_file)

    elif kind == "md":
        return partition_md(filename=temp_file)

    else:
        raise ValueError(f"Unsupported file_type: {file_type}")


# ─────────────────────────────────────────────
#  Element analysis & chunking helpers
# ─────────────────────────────────────────────

def analyze_elements(elements: list) -> dict:
    """Return a count breakdown of element types for UI display."""
    counts = {"text": 0, "tables": 0, "images": 0, "titles": 0, "other": 0}

    for element in elements:
        name = type(element).__name__
        if name == "Table":
            counts["tables"] += 1
        elif name == "Image":
            counts["images"] += 1
        elif name in ("Title", "Header"):
            counts["titles"] += 1
        elif name in ("NarrativeText", "Text", "ListItem", "FigureCaption"):
            counts["text"] += 1
        else:
            counts["other"] += 1

    logger.info("elements_analyzed", **counts)
    return counts


def separate_content_types(chunk, source_type: str = "file") -> dict:
    """
    Split a chunk's orig_elements into typed buckets:
    text, tables (HTML), images (base64), and a types list.
    """
    is_url_source = source_type == "url"
    content_data = {
        "text": chunk.text,
        "tables": [],
        "images": [],
        "types": ["text"],
    }

    orig = getattr(getattr(chunk, "metadata", None), "orig_elements", None)
    if not orig:
        return content_data

    for element in orig:
        element_type = type(element).__name__

        if element_type == "Table":
            content_data["types"].append("table")
            table_html = getattr(element.metadata, "text_as_html", element.text)
            content_data["tables"].append(table_html)

        elif element_type == "Image" and not is_url_source:
            b64 = getattr(getattr(element, "metadata", None), "image_base64", None)
            if b64:
                content_data["types"].append("image")
                content_data["images"].append(b64)

    content_data["types"] = list(set(content_data["types"]))
    return content_data


def get_page_number(chunk, chunk_index: int) -> int:
    """Return page number from chunk metadata, falling back to chunk index + 1."""
    page = getattr(getattr(chunk, "metadata", None), "page_number", None)
    return page if page is not None else chunk_index + 1


def create_ai_summary(text: str, tables_html: list, images_base64: list) -> str:
    """
    Generate a searchable index summary for a chunk that contains tables or images.
    Uses gpt-4o via LangChain so vision content (images) is supported.
    """
    try:
        prompt_text = f"""Create a searchable index for this document content.
CONTENT:
{text}
"""
        if tables_html:
            prompt_text += "TABLES:\n"
            for i, table in enumerate(tables_html):
                prompt_text += f"Table {i + 1}:\n{table}\n\n"

        prompt_text += """Generate a structured search index (aim for 250-400 words):

QUESTIONS: List 5-7 key questions this content answers (use what/how/why/when/who variations)

KEYWORDS: Include:
- Specific data (numbers, dates, percentages, amounts)
- Core concepts and themes
- Technical terms and casual alternatives
- Industry terminology

VISUALS (if images present):
- Chart/graph types and what they show
- Trends and patterns visible
- Key insights from visualizations

DATA RELATIONSHIPS (if tables present):
- Column headers and their meaning
- Key metrics and relationships
- Notable values or patterns

Focus on terms users would actually search for. Be specific and comprehensive.

SEARCH INDEX:"""

        message_content = [{"type": "text", "text": prompt_text}]
        for image_base64 in images_base64:
            message_content.append(
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
                }
            )

        message = HumanMessage(content=message_content)
        response = openAI["embeddings_llm"].invoke([message])
        return response.content

    except Exception as e:
        raise Exception(f"Failed to create AI summary: {str(e)}")